"""
CyberShield - Documentation Generator
Generates: Full Report PDF, Abstract PDF, PowerPoint, ZIP
All strings are Latin-1 safe for fpdf2 compatibility.
"""

import os, zipfile
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "docs_output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

CYAN   = RGBColor(0x00, 0xE6, 0xFF)
GREEN  = RGBColor(0x00, 0xFF, 0x88)
RED    = RGBColor(0xFF, 0x17, 0x44)
DARK   = RGBColor(0x0D, 0x11, 0x1F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0x8A, 0x9B, 0xBF)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)
PURPLE = RGBColor(0x7C, 0x4D, 0xFF)

# ==============================================================
#  PDF
# ==============================================================
class ReportPDF(FPDF):
    def header(self):
        self.set_fill_color(13, 17, 31)
        self.rect(0, 0, 210, 18, 'F')
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(0, 200, 230)
        self.set_y(5)
        self.cell(0, 8, "CyberShield - AI-Powered Cyber Threat Detection System", align="C")
        self.ln(14)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(138, 155, 191)
        self.cell(0, 10, f"Page {self.page_no()} | CyberShield Project Report | 2025-26", align="C")

    def chapter_title(self, title):
        self.set_font("Helvetica", "B", 14)
        self.set_fill_color(0, 200, 230)
        self.set_text_color(13, 17, 31)
        self.cell(0, 9, "  " + title, ln=True, fill=True)
        self.ln(4)
        self.set_text_color(30, 35, 55)

    def section_title(self, title):
        self.set_font("Helvetica", "B", 11)
        self.set_text_color(0, 160, 200)
        self.cell(0, 7, title, ln=True)
        self.set_draw_color(0, 160, 200)
        self.set_line_width(0.3)
        self.line(self.get_x(), self.get_y(), 200, self.get_y())
        self.ln(3)
        self.set_text_color(30, 35, 55)

    def body(self, text):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(40, 45, 65)
        self.multi_cell(0, 6, text)
        self.ln(3)

    def bullet(self, items):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(40, 45, 65)
        for item in items:
            self.set_x(15)
            self.multi_cell(0, 6, "  - " + item)

    def kv_table(self, rows, col_w=(70, 110)):
        self.set_font("Helvetica", "", 10)
        for i, (k, v) in enumerate(rows):
            fill = (i % 2 == 0)
            self.set_fill_color(235, 240, 252) if fill else self.set_fill_color(255, 255, 255)
            self.set_text_color(0, 120, 160)
            self.cell(col_w[0], 7, "  " + k, border=1, fill=fill)
            self.set_text_color(40, 45, 65)
            self.cell(col_w[1], 7, "  " + v, border=1, fill=fill, ln=True)
        self.ln(4)

    def metric_table(self, headers, rows):
        col_w = 190 // len(headers)
        self.set_font("Helvetica", "B", 10)
        self.set_fill_color(13, 17, 31)
        self.set_text_color(0, 200, 230)
        for h in headers:
            self.cell(col_w, 8, h, border=1, fill=True, align="C")
        self.ln()
        self.set_font("Helvetica", "", 10)
        for i, row in enumerate(rows):
            self.set_fill_color(235, 245, 252) if i % 2 == 0 else self.set_fill_color(255, 255, 255)
            self.set_text_color(30, 35, 55)
            for cell in row:
                self.cell(col_w, 7, str(cell), border=1, fill=True, align="C")
            self.ln()
        self.ln(4)


def build_pdf():
    pdf = ReportPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.set_margins(15, 22, 15)

    # ---- COVER ----
    pdf.add_page()
    pdf.set_fill_color(13, 17, 31)
    pdf.rect(0, 0, 210, 297, 'F')
    pdf.set_y(55)
    pdf.set_font("Helvetica", "B", 36)
    pdf.set_text_color(0, 200, 230)
    pdf.cell(0, 18, "CYBERSHIELD", align="C", ln=True)
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_text_color(0, 240, 130)
    pdf.cell(0, 10, "AI-Powered Cyber Threat Detection System", align="C", ln=True)
    pdf.ln(6)
    pdf.set_draw_color(0, 200, 230)
    pdf.set_line_width(0.6)
    pdf.line(30, pdf.get_y(), 180, pdf.get_y())
    pdf.ln(8)
    pdf.set_font("Helvetica", "", 12)
    pdf.set_text_color(138, 155, 191)
    for line in [
        "Complete Technical Report",
        "",
        "Final Year B.E./B.Tech Project -- 2025-26",
        "Department of Computer Science and Engineering",
        "",
        "Machine Learning  |  Flask REST API  |  Real-Time Dashboard",
        "Random Forest  |  XGBoost  |  Neural Network (MLP)",
        "",
        "Random Forest: 99.9%   XGBoost: 99.87%   MLP: 99.64%",
    ]:
        pdf.cell(0, 8, line, align="C", ln=True)

    # ---- TOC ----
    pdf.add_page()
    pdf.chapter_title("Table of Contents")
    toc = [
        ("1.", "Project Overview and Objectives"),
        ("2.", "System Architecture"),
        ("3.", "Technology Stack"),
        ("4.", "Backend - Flask REST API (Deep Dive)"),
        ("5.", "Machine Learning Pipeline"),
        ("6.", "Frontend - Dashboard (Deep Dive)"),
        ("7.", "Data Flow - How It All Works Together"),
        ("8.", "Attack Types Detected"),
        ("9.", "Model Performance Metrics"),
        ("10.", "File Structure"),
        ("11.", "Setup and Deployment Guide"),
        ("12.", "Limitations and Future Work"),
    ]
    pdf.set_font("Helvetica", "", 11)
    for num, title in toc:
        pdf.set_text_color(0, 120, 160)
        pdf.cell(14, 8, num)
        pdf.set_text_color(30, 35, 55)
        pdf.cell(0, 8, title, ln=True)
        pdf.set_draw_color(200, 210, 230)
        pdf.set_line_width(0.2)
        pdf.line(15, pdf.get_y(), 195, pdf.get_y())

    # ---- SECTION 1 ----
    pdf.add_page()
    pdf.chapter_title("1. Project Overview and Objectives")
    pdf.body(
        "CyberShield is a final year engineering project that demonstrates a complete, "
        "end-to-end AI-powered Cyber Threat Detection System. The system is composed of "
        "two major components: a Python-based Machine Learning backend (served via a Flask "
        "REST API) and a cyberpunk-themed real-time web dashboard built in HTML, CSS, and "
        "JavaScript.\n\n"
        "The system simulates the detection of seven categories of network threats using "
        "three state-of-the-art machine learning models trained on a synthetic dataset of "
        "100,000 network traffic records. The dashboard visualises threat feeds, attack "
        "distributions, model performance, and supports manual packet analysis in real time."
    )
    pdf.section_title("Primary Objectives")
    pdf.bullet([
        "Train multiple ML models to classify network traffic as normal or as one of six attack types.",
        "Deploy a Flask REST API that serves real-time predictions and simulation data.",
        "Build an interactive web dashboard that visualises live threat intelligence.",
        "Demonstrate ensemble ML decision-making using three independent models.",
        "Provide a manual packet analyser for interactive demonstration to examiners.",
    ])
    pdf.section_title("Key Highlights")
    pdf.kv_table([
        ("Project Name",         "CyberShield - AI Cyber Threat Detection"),
        ("Domain",               "Cybersecurity + Machine Learning"),
        ("Language",             "Python 3.x (Backend), HTML/CSS/JS (Frontend)"),
        ("ML Models",            "Random Forest, XGBoost, Neural Network (MLP)"),
        ("Dataset Size",         "100,000 synthetic network traffic records"),
        ("Best Model Accuracy",  "99.9% (Random Forest)"),
        ("API Framework",        "Flask + Flask-CORS"),
        ("Frontend Framework",   "Vanilla HTML5, CSS3, JavaScript ES6+"),
        ("Charts Library",       "Chart.js v4.4.0"),
        ("Attack Types",         "DDoS, Port Scan, Brute Force, SQL Injection, XSS, Botnet"),
    ])

    # ---- SECTION 2 ----
    pdf.add_page()
    pdf.chapter_title("2. System Architecture")
    pdf.body(
        "CyberShield follows a classic two-tier Client-Server architecture, with a clear "
        "separation between the frontend presentation layer and the backend intelligence layer.\n\n"
        "The browser (Client) loads the dashboard once and establishes a polling connection "
        "to the Flask API Server (Server) running locally on port 5000. The server loads "
        "pre-trained ML models from disk and responds to API requests with JSON data."
    )
    pdf.section_title("Architecture Diagram (Text)")
    pdf.set_font("Courier", "", 9)
    pdf.set_fill_color(240, 244, 255)
    pdf.set_text_color(30, 35, 55)
    arch = (
        "  +----------------------------+      HTTP / JSON\n"
        "  |  FRONTEND (Browser)        | <--------------> +-------------------------+\n"
        "  |  index.html                |                  |  BACKEND (Flask API)    |\n"
        "  |  css/style.css             | GET /api/simulate|  app.py  port:5000      |\n"
        "  |  js/main.js                | GET /api/status  |                         |\n"
        "  |                            | GET /api/stats   |  +-------------------+  |\n"
        "  |  - Live Threat Feed        | GET /api/history |  |  ML Models (.pkl) |  |\n"
        "  |  - Traffic Timeline Chart  | POST /api/predict|  |  random_forest    |  |\n"
        "  |  - Threat Dist Donut       | POST /api/reset  |  |  xgboost          |  |\n"
        "  |  - Model Performance Bars  |                  |  |  mlp              |  |\n"
        "  |  - Packet Analyser Form    |                  |  +-------------------+  |\n"
        "  |  - Threat Log Table        |                  |  data_generator.py      |\n"
        "  +----------------------------+                  |  threat_detector.py     |\n"
        "                                                  +-------------------------+\n"
    )
    pdf.multi_cell(0, 5, arch, fill=True)
    pdf.ln(3)
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(40, 45, 65)

    # ---- SECTION 3 ----
    pdf.add_page()
    pdf.chapter_title("3. Technology Stack")
    pdf.section_title("Backend Technologies")
    pdf.kv_table([
        ("Python 3.x",      "Core programming language for all backend logic"),
        ("Flask 3.x",       "Lightweight WSGI web framework for the REST API"),
        ("Flask-CORS",      "Cross-Origin Resource Sharing for browser access"),
        ("scikit-learn",    "Random Forest classifier + StandardScaler + LabelEncoder"),
        ("XGBoost",         "Gradient-boosted decision tree classifier"),
        ("pandas",          "Data manipulation and CSV handling"),
        ("NumPy",           "Numerical computing for feature arrays"),
        ("joblib/pickle",   "Model serialisation and deserialisation (.pkl files)"),
    ])
    pdf.section_title("Frontend Technologies")
    pdf.kv_table([
        ("HTML5",           "Semantic markup structure of the dashboard"),
        ("CSS3",            "Custom cyberpunk styling with CSS variables and animations"),
        ("JavaScript ES6+", "Async/Await fetch API, DOM manipulation, event handling"),
        ("Chart.js v4.4.0", "Line chart (Traffic Timeline) and Doughnut chart (Distribution)"),
        ("Google Fonts",    "Orbitron, Rajdhani, Share Tech Mono - cyberpunk typefaces"),
        ("Canvas API",      "Animated floating particle background effect"),
    ])
    pdf.section_title("Data and Models")
    pdf.kv_table([
        ("network_traffic.csv", "100,000 synthetic records with 20 network features"),
        ("random_forest.pkl",   "Trained Random Forest model (~4.2 MB)"),
        ("xgboost.pkl",         "Trained XGBoost model (~1.2 MB)"),
        ("mlp.pkl",             "Trained Neural Network model (~1.1 MB)"),
        ("scaler.pkl",          "StandardScaler fitted on training data"),
        ("label_encoder.pkl",   "Label encoder mapping class strings to integers"),
        ("metrics.json",        "Stored accuracy/F1/precision/recall for all models"),
    ])

    # ---- SECTION 4 ----
    pdf.add_page()
    pdf.chapter_title("4. Backend - Flask REST API (Deep Dive)")
    pdf.body(
        "The backend is the intelligence layer of CyberShield. It is a Python application "
        "composed of four files: app.py (API routes), train_model.py (model training), "
        "data_generator.py (synthetic data generation), and threat_detector.py (prediction engine)."
    )
    pdf.section_title("4.1  data_generator.py - Synthetic Dataset Generator")
    pdf.body(
        "This module generates a synthetic dataset of 100,000 network traffic records. "
        "Each record contains 20 handcrafted numerical features that mimic real KDD Cup "
        "network intrusion data. Each attack class has a unique statistical signature:"
    )
    pdf.bullet([
        "DDoS: Very high connection count (400-500), high SYN error rate (0.8-1.0), ICMP protocol",
        "Port Scan: High count, zero destination bytes, high SYN error rate, TCP protocol",
        "Brute Force: High failed login attempts (10-20), elevated SYN errors, moderate bytes",
        "SQL Injection: High hot indicators (15-30), elevated source/destination bytes, long duration",
        "XSS: Moderate hot indicators, elevated bytes, shorter duration than SQL injection",
        "Botnet: Very long duration (100-300s), UDP protocol, steady byte flow, low error rates",
        "Normal: Low counts, low error rates, balanced bytes, typical session durations",
    ])
    pdf.section_title("4.2  train_model.py - ML Training Pipeline")
    pdf.body(
        "This script trains all three ML models in sequence:\n"
        "  Step 1: Load network_traffic.csv using pandas\n"
        "  Step 2: Scale features using StandardScaler (zero mean, unit variance)\n"
        "  Step 3: Encode labels (attack types) using LabelEncoder\n"
        "  Step 4: Split data - 80% training, 20% testing (stratified)\n"
        "  Step 5: Train Random Forest (200 estimators, max_depth=20)\n"
        "  Step 6: Train XGBoost (300 estimators, max_depth=8, eta=0.1)\n"
        "  Step 7: Train MLP Neural Network (3 hidden layers: 256-128-64 neurons)\n"
        "  Step 8: Save all models and metrics to models/ directory"
    )
    pdf.section_title("4.3  threat_detector.py - Prediction Engine")
    pdf.body(
        "The ThreatDetector class loads all three trained models and exposes two methods:\n\n"
        "predict(features, model_name): Runs a single model on input features and returns "
        "threat_type, severity, confidence, description, recommendation, and probabilities.\n\n"
        "predict_ensemble(features): Runs all three models, averages their probability "
        "outputs, and returns a combined prediction - making the classification more robust "
        "than any single model alone.\n\n"
        "Severity mapping:\n"
        "  critical: DDoS, Botnet\n"
        "  high: Brute Force, SQL Injection\n"
        "  medium: XSS, Port Scan\n"
        "  none: Normal traffic"
    )
    pdf.section_title("4.4  app.py - REST API Endpoints")
    pdf.metric_table(
        ["Method", "Endpoint", "Description"],
        [
            ["GET",  "/api/status",             "System health + model status + metrics"],
            ["POST", "/api/predict",            "Classify a manually submitted packet"],
            ["GET",  "/api/simulate",           "Generate and classify a random packet"],
            ["GET",  "/api/history",            "Return recent threat log (last 40)"],
            ["GET",  "/api/stats",              "Aggregate counters (total, critical...)"],
            ["GET",  "/api/threat-distribution","Count of each attack type seen"],
            ["POST", "/api/reset",              "Clear in-memory state"],
        ]
    )

    # ---- SECTION 5 ----
    pdf.add_page()
    pdf.chapter_title("5. Machine Learning Pipeline")
    pdf.section_title("Feature Set (20 Features)")
    pdf.kv_table([
        ("duration",                 "Length of the connection in seconds"),
        ("protocol_type",            "Network protocol: 0=TCP, 1=UDP, 2=ICMP"),
        ("src_bytes",                "Bytes sent from source to destination"),
        ("dst_bytes",                "Bytes sent from destination to source"),
        ("wrong_fragment",           "Number of wrong fragments in packet"),
        ("urgent",                   "Number of urgent packets"),
        ("hot",                      "Hot indicators (suspicious access patterns)"),
        ("num_failed_logins",        "Number of failed login attempts"),
        ("logged_in",                "1 if successfully logged in, 0 otherwise"),
        ("num_compromised",          "Number of compromised conditions"),
        ("count",                    "Connections to same host in last 2 seconds"),
        ("srv_count",                "Connections to same service in last 2 seconds"),
        ("serror_rate",              "Fraction of connections with SYN errors"),
        ("rerror_rate",              "Fraction of connections with REJ errors"),
        ("same_srv_rate",            "Fraction of connections to the same service"),
        ("diff_srv_rate",            "Fraction of connections to different services"),
        ("dst_host_count",           "Count of connections to destination host"),
        ("dst_host_srv_count",       "Count of connections with same service on dst"),
        ("dst_host_same_srv_rate",   "Fraction same service on destination host"),
        ("dst_host_serror_rate",     "Fraction of SYN errors on destination host"),
    ], col_w=(65, 115))
    pdf.section_title("Model Architectures")
    pdf.bullet([
        "Random Forest: Ensemble of 200 decision trees, each trained on a random feature subset. "
        "Final prediction is by majority voting. Handles non-linear boundaries well.",
        "XGBoost: Gradient-boosted tree ensemble with 300 estimators. Builds trees sequentially, "
        "each correcting errors of the previous. Uses learning rate=0.1 and max_depth=8.",
        "MLP Neural Network: 3-layer feed-forward network. Input (20) -> Hidden1 (256, ReLU) -> "
        "Hidden2 (128, ReLU) -> Hidden3 (64, ReLU) -> Output (7 classes, Softmax). Adam optimiser.",
    ])

    # ---- SECTION 6 ----
    pdf.add_page()
    pdf.chapter_title("6. Frontend - Dashboard (Deep Dive)")
    pdf.body(
        "The frontend is a single-page application (SPA) built entirely with vanilla HTML5, "
        "CSS3, and ES6+ JavaScript. No JavaScript framework (React, Vue, Angular) is used - "
        "this demonstrates strong fundamentals of pure web development.\n\n"
        "The UI is themed in a cyberpunk aesthetic using dark backgrounds, neon cyan/green "
        "accents, scanline effects, and animated particles."
    )
    pdf.section_title("6.1  Structure - index.html")
    pdf.bullet([
        "Animated Background: Full-page Canvas particle system + CSS grid overlay + scanline shimmer",
        "Header: Logo, SYSTEM ONLINE/OFFLINE status indicator with pulsing dot, 24h live clock, "
        "LIVE SCAN toggle button, RESET button",
        "Alert Banner: Slides down with colour-coded alerts for HIGH/CRITICAL threats",
        "Stats Row: 6 stat cards - Packets Analysed, Threats Detected, Critical, High, Medium, Low",
        "Left Column: Live Threat Feed (scrolling real-time log) + Traffic Timeline Line Chart",
        "Middle Column: Threat Distribution Donut Chart + Model Performance Bars + Attack Reference",
        "Right Column: Interactive Packet Analyser Form + Threat Log Table",
    ])
    pdf.section_title("6.2  Styling - css/style.css")
    pdf.bullet([
        "CSS Custom Properties: All colours/sizes defined as --var for easy theming",
        "Glassmorphism panels: backdrop-filter blur, semi-transparent backgrounds",
        "Animations: keyframes for scan pulses, glow flickers, feed item slide-in",
        "Responsive Grid: CSS Grid 3-column layout, collapses to single column on mobile",
        "Status Dot: CSS animation alternates between pulsing green (online) and red (offline)",
        "Hover Effects: All interactive cards and buttons have transform + box-shadow transitions",
    ])
    pdf.section_title("6.3  Logic - js/main.js")
    pdf.bullet([
        "boot(): Entry point on DOMContentLoaded - starts clock, builds legend, fetches status, "
        "initialises charts, registers listeners, auto-starts the live scan",
        "loadStatus(): Calls GET /api/status - sets ONLINE/OFFLINE indicator, loads model metrics",
        "startScan() / runSimulate(): Polls GET /api/simulate every 1.8 seconds, feeds results into "
        "the threat feed, chart buffers, and stat cards",
        "analyzePacket(): Reads form inputs, POSTs to /api/predict, displays colour-coded result "
        "with probability bars for all 7 classes",
        "initCharts(): Creates Chart.js Line chart (Traffic Timeline) and Doughnut chart "
        "(Threat Distribution) with custom cyberpunk colour schemes",
        "applyPreset(): Pre-fills the packet analyser form with realistic values for each attack type",
        "refreshLog(): Periodically fetches GET /api/history and renders tabular threat log",
        "initParticles(): Canvas 2D floating particle system with connecting lines, runs on RAF loop",
    ])

    # ---- SECTION 7 ----
    pdf.add_page()
    pdf.chapter_title("7. Data Flow - How It All Works Together")
    pdf.set_font("Courier", "", 9)
    pdf.set_fill_color(240, 244, 255)
    pdf.set_text_color(30, 35, 55)
    flow = (
        "USER opens index.html in browser\n"
        "        |\n"
        "        v\n"
        "js/main.js boot() runs\n"
        "    |- Fetches GET /api/status -> Flask checks models -> returns JSON\n"
        "    |- Sets header ONLINE/OFFLINE indicator\n"
        "    |- Loads model accuracy bars from metrics in status response\n"
        "        |\n"
        "        v\n"
        "startScan() starts setInterval every 1800ms\n"
        "        |\n"
        "        v\n"
        "runSimulate() fires:\n"
        "    |- GET /api/simulate\n"
        "          |\n"
        "          v\n"
        "    Flask simulate():\n"
        "        |- Randomly picks attack type (weighted)\n"
        "        |- Calls data_generator._generate_features(type)\n"
        "        |- Calls ThreatDetector.predict_ensemble(features)\n"
        "              |- Runs Random Forest  -> probability vector\n"
        "              |- Runs XGBoost        -> probability vector\n"
        "              |- Runs MLP            -> probability vector\n"
        "              |- Averages 3 vectors  -> final prediction\n"
        "        |- Returns {threat_type, severity, confidence, src_ip, ...}\n"
        "          |\n"
        "          v\n"
        "    JS receives response:\n"
        "        |- addFeedItem()   -> scroll new row into Live Feed\n"
        "        |- distCounts[]++ -> update Donut chart slice\n"
        "        |- fetchStats()   -> update 6 stat counter cards\n"
        "        |- showAlert()    -> slide-down banner if HIGH/CRITICAL\n"
        "        |- refreshLog()   -> update Threat Log table every 3 calls\n"
        "        |\n"
        "        v\n"
        "[Repeat every 1.8 seconds]\n\n"
        "MANUAL PACKET ANALYSIS:\n"
        "User fills form -> analyzePacket()\n"
        "    |- POST /api/predict {features, model, ensemble}\n"
        "    |- Flask predict() -> ThreatDetector.predict()\n"
        "    |- Returns full result with probabilities\n"
        "    |- showAnalysisResult() renders result + probability bars\n"
    )
    pdf.multi_cell(0, 5, flow, fill=True)
    pdf.ln(2)
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(40, 45, 65)

    # ---- SECTION 8 ----
    pdf.add_page()
    pdf.chapter_title("8. Attack Types Detected")
    attacks = [
        ("Normal Traffic", "none",     (0, 160, 80),
         "Regular benign network communication. Low counts, balanced bytes, no error spikes."),
        ("DDoS",           "critical", (200, 20, 50),
         "Distributed Denial of Service: floods the server with ICMP packets at high rate. "
         "Characterised by extreme connection counts (400-500) and SYN error rates near 1.0."),
        ("Port Scan",      "medium",   (180, 120, 0),
         "Attacker probes target for open ports using TCP SYN packets. "
         "Sends packets to hundreds of ports with zero response bytes expected."),
        ("Brute Force",    "high",     (200, 60, 0),
         "Repeated login attempts guessing credentials. Identified by high "
         "num_failed_logins (10-20) and elevated SYN error rates."),
        ("SQL Injection",  "high",     (200, 60, 0),
         "Malicious SQL code injected into web form inputs. "
         "Detected by high hot indicator scores and large byte transfers."),
        ("XSS",            "medium",   (180, 120, 0),
         "Cross-Site Scripting: attacker injects client-side scripts into web pages. "
         "Moderate hot indicators, shorter sessions than SQL injection."),
        ("Botnet C&C",     "critical", (200, 20, 50),
         "Compromised machine communicating with Command & Control server. "
         "Characterised by long session duration (100-300s), UDP protocol, steady traffic."),
    ]
    for name, sev, col, desc in attacks:
        pdf.set_font("Helvetica", "B", 11)
        pdf.set_text_color(*col)
        pdf.cell(0, 8, f"{name}  [Severity: {sev.upper()}]", ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(40, 45, 65)
        pdf.multi_cell(0, 6, desc)
        pdf.ln(2)

    # ---- SECTION 9 ----
    pdf.add_page()
    pdf.chapter_title("9. Model Performance Metrics")
    pdf.body(
        "All three models were evaluated on a held-out 20% test set (20,000 records) "
        "using four standard classification metrics:"
    )
    pdf.metric_table(
        ["Model", "Accuracy", "Precision", "Recall", "F1-Score"],
        [
            ["Random Forest",        "99.90%", "99.90%", "99.90%", "99.90%"],
            ["XGBoost",              "99.87%", "99.87%", "99.87%", "99.87%"],
            ["Neural Network (MLP)", "99.64%", "99.65%", "99.64%", "99.64%"],
        ]
    )
    pdf.section_title("Metric Definitions")
    pdf.kv_table([
        ("Accuracy",  "% of total predictions that were correct"),
        ("Precision", "Of all packets labelled as attack, % that were actually attacks"),
        ("Recall",    "Of all actual attacks, % that were correctly caught"),
        ("F1-Score",  "Harmonic mean of Precision and Recall - balanced single metric"),
    ])
    pdf.section_title("Why These Scores Are High")
    pdf.body(
        "The models achieve near-perfect accuracy because the synthetic dataset has clearly "
        "differentiated statistical signatures for each attack class (e.g. DDoS always has "
        "count > 400, while normal traffic has count < 50). In a real-world deployment with "
        "noisy, overlapping network traffic, accuracy would be lower (typically 92-97% for "
        "similar models on the KDD Cup 1999 dataset). These results are valid and expected "
        "for a clean synthetic demonstration dataset."
    )

    # ---- SECTION 10 ----
    pdf.add_page()
    pdf.chapter_title("10. File Structure")
    pdf.set_font("Courier", "", 9)
    pdf.set_fill_color(240, 244, 250)
    pdf.set_text_color(30, 35, 55)
    tree = (
        "CyberShield/\n"
        "|\n"
        "+-- README.md                   Project overview and setup guide\n"
        "+-- setup_and_run.bat           One-click Windows setup + launch script\n"
        "+-- start_server.bat            Quick Flask server launcher\n"
        "|\n"
        "+-- backend/\n"
        "|   +-- app.py                  Flask REST API (7 endpoints)\n"
        "|   +-- train_model.py          ML training pipeline\n"
        "|   +-- data_generator.py       Synthetic dataset + simulation engine\n"
        "|   +-- threat_detector.py      ML prediction + ensemble logic\n"
        "|   +-- requirements.txt        Python dependencies\n"
        "|   |\n"
        "|   +-- data/\n"
        "|   |   +-- network_traffic.csv 100,000 synthetic records (~10MB)\n"
        "|   |\n"
        "|   +-- models/\n"
        "|       +-- random_forest.pkl   Trained Random Forest (~4.2MB)\n"
        "|       +-- xgboost.pkl         Trained XGBoost (~1.2MB)\n"
        "|       +-- mlp.pkl             Trained MLP Neural Network (~1.1MB)\n"
        "|       +-- scaler.pkl          StandardScaler fitted on training data\n"
        "|       +-- label_encoder.pkl   Attack type label encoder\n"
        "|       +-- metrics.json        Stored evaluation metrics\n"
        "|\n"
        "+-- frontend/\n"
        "    +-- index.html              Single-page dashboard (290 lines)\n"
        "    +-- css/\n"
        "    |   +-- style.css           Cyberpunk theme styles\n"
        "    +-- js/\n"
        "        +-- main.js             Dashboard controller logic (~595 lines)\n"
    )
    pdf.multi_cell(0, 5, tree, fill=True)
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(40, 45, 65)

    # ---- SECTION 11 ----
    pdf.add_page()
    pdf.chapter_title("11. Setup and Deployment Guide")
    pdf.section_title("Prerequisites")
    pdf.bullet([
        "Python 3.8 or higher installed",
        "pip package manager",
        "Modern web browser (Chrome, Edge, Firefox)",
    ])
    pdf.section_title("Installation Steps")
    pdf.set_font("Courier", "", 9)
    pdf.set_fill_color(240, 244, 250)
    pdf.set_text_color(30, 35, 55)
    steps = (
        "# Step 1: Navigate to backend directory\n"
        "cd CyberShield/backend\n\n"
        "# Step 2: Install Python dependencies\n"
        "pip install flask flask-cors scikit-learn xgboost pandas numpy joblib\n\n"
        "# Step 3: Generate dataset and train models (one time only)\n"
        "python train_model.py\n\n"
        "# Step 4: Start the Flask API server\n"
        "python app.py\n"
        "# Server starts at: http://localhost:5000\n\n"
        "# Step 5: Open the dashboard\n"
        "# Open frontend/index.html in any web browser\n"
        "# Status indicator should show: SYSTEM ONLINE\n"
    )
    pdf.multi_cell(0, 5, steps, fill=True)
    pdf.ln(3)
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(40, 45, 65)
    pdf.body(
        "Alternatively, double-click setup_and_run.bat on Windows to automatically "
        "install dependencies, train models, start the server, and open the browser."
    )

    # ---- SECTION 12 ----
    pdf.add_page()
    pdf.chapter_title("12. Limitations and Future Work")
    pdf.section_title("Current Limitations")
    pdf.bullet([
        "Synthetic Data: The system uses generated data, not real network captures.",
        "Local Only: Flask server runs on localhost; production needs Gunicorn + HTTPS.",
        "In-Memory State: Threat history is stored in RAM; restarting clears all logs.",
        "No Authentication: The dashboard and API have no login protection.",
        "Single Machine: Does not monitor distributed/cloud network infrastructure.",
    ])
    pdf.section_title("Future Enhancements")
    pdf.bullet([
        "Real-time Packet Capture: Integrate Scapy or libpcap to analyse actual traffic.",
        "Persistent Storage: Use SQLite or PostgreSQL to store threat logs.",
        "Deep Learning: Implement LSTM/Autoencoder models for time-series anomaly detection.",
        "Alerting System: Email/SMS notifications via Twilio/SendGrid.",
        "Docker Deployment: Containerise backend and frontend for cloud deployment.",
        "User Authentication: Add JWT-based login for the dashboard.",
        "Threat Intelligence: Integrate with VirusTotal or AbuseIPDB APIs.",
    ])

    out_path = os.path.join(OUTPUT_DIR, "CyberShield_Full_Report.pdf")
    pdf.output(out_path)
    print("[OK] Full Report PDF saved: " + out_path)
    return out_path


# ==============================================================
#  ABSTRACT PDF
# ==============================================================
def build_abstract():
    pdf = FPDF()
    pdf.add_page()
    pdf.set_margins(20, 20, 20)
    pdf.set_auto_page_break(auto=True, margin=20)

    pdf.set_fill_color(13, 17, 31)
    pdf.rect(0, 0, 210, 40, 'F')
    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(0, 200, 230)
    pdf.set_y(8)
    pdf.cell(0, 10, "CYBERSHIELD", align="C", ln=True)
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(0, 240, 130)
    pdf.cell(0, 8, "AI-Powered Cyber Threat Detection System", align="C", ln=True)
    pdf.ln(14)

    pdf.set_font("Helvetica", "B", 13)
    pdf.set_text_color(0, 120, 160)
    pdf.cell(0, 9, "ABSTRACT", align="C", ln=True)
    pdf.set_draw_color(0, 160, 210)
    pdf.set_line_width(0.5)
    pdf.line(20, pdf.get_y(), 190, pdf.get_y())
    pdf.ln(6)

    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(30, 35, 55)
    abstract_text = (
        "In today's rapidly evolving digital landscape, cyber threats such as Distributed Denial "
        "of Service (DDoS) attacks, brute force intrusions, SQL injections, and botnet communications "
        "pose severe risks to enterprise networks and critical infrastructure. Traditional rule-based "
        "intrusion detection systems (IDS) are increasingly inadequate against sophisticated, "
        "adaptive attack vectors. This project presents CyberShield, a comprehensive, AI-powered "
        "Cyber Threat Detection System that leverages ensemble machine learning to classify network "
        "traffic in real time.\n\n"
        "CyberShield is implemented as a full-stack web application with two core components: "
        "a Python Flask REST API backend and an interactive cyberpunk-themed web dashboard frontend. "
        "The backend houses three independently trained machine learning models - Random Forest, "
        "XGBoost, and a Multi-Layer Perceptron (MLP) Neural Network - each trained on a synthetic "
        "dataset of 100,000 network traffic records containing 20 carefully engineered network "
        "features derived from the KDD Cup intrusion detection benchmark.\n\n"
        "The ensemble prediction mechanism averages the probability outputs of all three models, "
        "reducing individual model bias and improving overall classification robustness. The system "
        "detects seven traffic categories: normal traffic, DDoS, port scanning, brute force, SQL "
        "injection, cross-site scripting (XSS), and botnet command-and-control communication. "
        "Each prediction includes a threat severity rating (none/low/medium/high/critical), "
        "a confidence percentage, and actionable security recommendations.\n\n"
        "The frontend dashboard, developed in vanilla HTML5, CSS3, and JavaScript, visualises "
        "threat intelligence in real time using Chart.js for an animated traffic timeline and "
        "threat distribution chart. The dashboard polls the /api/simulate endpoint every 1.8 "
        "seconds to generate and classify synthetic network packets, simulating a live threat "
        "monitoring environment. A dedicated Packet Analyser panel allows users to manually "
        "submit packet feature values using predefined attack presets (DDoS, Port Scan, Brute "
        "Force, SQL Injection, Botnet) and receive instant multi-model analysis results.\n\n"
        "Experimental evaluation on the 20,000-record test set demonstrates exceptional "
        "performance: Random Forest achieves 99.90% accuracy, XGBoost achieves 99.87%, and "
        "the MLP Neural Network achieves 99.64% accuracy, with correspondingly high precision, "
        "recall, and F1-scores for all models. These results validate the effectiveness of the "
        "ensemble approach and the discriminative power of the engineered feature set.\n\n"
        "CyberShield demonstrates practical integration of machine learning, REST API design, "
        "and real-time web visualisation in the domain of cybersecurity. While the current "
        "implementation uses synthetic data for demonstration, the modular architecture is "
        "designed for extension to real packet capture using libraries such as Scapy or libpcap. "
        "Future work includes persistent threat logging, email alerting, deep learning (LSTM) "
        "models for sequential traffic analysis, and cloud deployment using Docker."
    )
    pdf.multi_cell(0, 7, abstract_text)
    pdf.ln(8)

    pdf.set_font("Helvetica", "B", 10)
    pdf.set_text_color(0, 120, 160)
    pdf.cell(35, 7, "Keywords:")
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(40, 45, 65)
    pdf.multi_cell(0, 7,
        "Intrusion Detection System, Machine Learning, Random Forest, XGBoost, Neural Network, "
        "Cybersecurity, Network Traffic Classification, Ensemble Learning, Flask REST API, "
        "Real-Time Dashboard, DDoS Detection, Anomaly Detection")

    out_path = os.path.join(OUTPUT_DIR, "CyberShield_Abstract.pdf")
    pdf.output(out_path)
    print("[OK] Abstract PDF saved: " + out_path)
    return out_path


# ==============================================================
#  POWERPOINT
# ==============================================================
def set_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_textbox(slide, text, left, top, width, height,
                size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_rect(slide, left, top, width, height, r, g, b):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    return shape

def add_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def bullet_slide(prs, title, points, title_color=CYAN):
    slide = add_slide(prs)
    set_bg(slide, 13, 17, 31)
    add_rect(slide, 0, 0, 10, 0.08, 0, 200, 230)
    add_textbox(slide, title, 0.3, 0.15, 9, 0.6, size=22, bold=True, color=title_color)
    add_rect(slide, 0.3, 0.82, 8.5, 0.03, 0, 160, 200)
    y = 1.0
    for pt in points:
        add_textbox(slide, ">  " + pt, 0.3, y, 9.2, 0.45, size=13, color=WHITE)
        y += 0.48
    return slide


def build_ppt():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: TITLE
    slide = add_slide(prs)
    set_bg(slide, 13, 17, 31)
    add_rect(slide, 0, 0, 10, 0.08, 0, 200, 230)
    add_rect(slide, 0, 7.42, 10, 0.08, 0, 240, 130)
    add_rect(slide, 0, 2.8, 10, 0.04, 0, 200, 230)
    add_textbox(slide, "CYBERSHIELD", 0, 0.9, 10, 1.4,
                size=58, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_textbox(slide, "AI-Powered Cyber Threat Detection System", 0, 2.25, 10, 0.6,
                size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Final Year B.E./B.Tech Project  |  2025-26", 0, 2.95, 10, 0.5,
                size=14, color=LGREY, align=PP_ALIGN.CENTER)
    add_textbox(slide,
        "Machine Learning  |  Flask REST API  |  Real-Time Dashboard\n"
        "Random Forest  |  XGBoost  |  Neural Network (MLP)",
        0, 3.55, 10, 0.8, size=13, color=LGREY, align=PP_ALIGN.CENTER)
    add_textbox(slide,
        "Best Accuracy: 99.9%  |  Dataset: 100,000 records  |  7 Attack Classes",
        0, 4.5, 10, 0.5, size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    # SLIDE 2: PROBLEM STATEMENT
    bullet_slide(prs, "Problem Statement", [
        "Cyber attacks are growing in frequency, sophistication, and damage",
        "Traditional rule-based IDS cannot adapt to new and evolving threats",
        "Manual monitoring of network logs is slow, error-prone, and unscalable",
        "Organisations need real-time, intelligent, automated threat classification",
        "Most tools are black-boxes - examiners and analysts cannot inspect decisions",
        "CyberShield solves this: explainable AI + real-time visual monitoring",
    ], title_color=RED)

    # SLIDE 3: SYSTEM OVERVIEW
    slide = add_slide(prs)
    set_bg(slide, 13, 17, 31)
    add_rect(slide, 0, 0, 10, 0.08, 0, 200, 230)
    add_textbox(slide, "System Overview", 0.3, 0.15, 9, 0.6, size=22, bold=True, color=CYAN)
    add_rect(slide, 0.3, 0.82, 8.5, 0.03, 0, 160, 200)
    add_rect(slide, 0.3, 1.1, 3.5, 5.5, 0, 40, 70)
    add_textbox(slide, "FRONTEND", 0.5, 1.2, 3.1, 0.4, size=14, bold=True, color=CYAN)
    add_textbox(slide,
        "index.html\nstyle.css\nmain.js\n\nChart.js Charts\nParticle Canvas\nPacket Analyser\nThreat Feed\nLive Clock",
        0.5, 1.65, 3.1, 4.5, size=11, color=LGREY)
    add_textbox(slide, "HTTP\n<->\nJSON", 3.9, 2.9, 1.5, 1.2,
                size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_rect(slide, 5.5, 1.1, 4.1, 5.5, 0, 40, 70)
    add_textbox(slide, "BACKEND (Flask :5000)", 5.65, 1.2, 3.8, 0.4, size=13, bold=True, color=GREEN)
    add_textbox(slide,
        "app.py  (7 REST endpoints)\ndata_generator.py\nthreat_detector.py\ntrain_model.py\n\n"
        "ML Models:\n  random_forest.pkl\n  xgboost.pkl\n  mlp.pkl\n\nmodels/metrics.json\ndata/network_traffic.csv",
        5.65, 1.65, 3.8, 4.5, size=11, color=LGREY)

    # SLIDE 4: TECH STACK
    slide = add_slide(prs)
    set_bg(slide, 13, 17, 31)
    add_rect(slide, 0, 0, 10, 0.08, 0, 200, 230)
    add_textbox(slide, "Technology Stack", 0.3, 0.15, 9, 0.6, size=22, bold=True, color=CYAN)
    add_rect(slide, 0.3, 0.82, 8.5, 0.03, 0, 160, 200)
    add_rect(slide, 0.3, 1.1, 4.3, 5.6, 0, 30, 60)
    add_textbox(slide, "BACKEND", 0.5, 1.2, 3.8, 0.4, size=14, bold=True, color=GREEN)
    be_items = [
        "Python 3.x - Core language",
        "Flask - REST API framework",
        "Flask-CORS - Cross-origin support",
        "scikit-learn - RF + Scaler + Encoder",
        "XGBoost - Gradient boosted trees",
        "pandas - Data manipulation",
        "NumPy - Numerical arrays",
        "joblib - Model serialisation",
    ]
    y = 1.65
    for item in be_items:
        add_textbox(slide, "> " + item, 0.5, y, 3.9, 0.38, size=11, color=WHITE)
        y += 0.4
    add_rect(slide, 5.0, 1.1, 4.6, 5.6, 0, 30, 60)
    add_textbox(slide, "FRONTEND", 5.2, 1.2, 4.1, 0.4, size=14, bold=True, color=CYAN)
    fe_items = [
        "HTML5 - Semantic structure",
        "CSS3 - Glassmorphism + animations",
        "JavaScript ES6+ - Async fetch/DOM",
        "Chart.js v4.4 - Interactive charts",
        "Canvas API - Particle background",
        "Google Fonts - Cyberpunk typography",
        "CSS Grid - Responsive 3-col layout",
        "Fetch API - Backend communication",
    ]
    y = 1.65
    for item in fe_items:
        add_textbox(slide, "> " + item, 5.2, y, 4.2, 0.38, size=11, color=WHITE)
        y += 0.4

    # SLIDE 5: ML PIPELINE
    bullet_slide(prs, "Machine Learning Pipeline", [
        "Dataset: 100,000 synthetic records, 20 features, 7 classes",
        "Preprocessing: StandardScaler (zero mean, unit variance) + LabelEncoder",
        "Train/Test Split: 80% training (80,000) / 20% testing (20,000) - stratified",
        "Random Forest: 200 decision trees, majority vote, max_depth=20",
        "XGBoost: 300 gradient-boosted trees, learning rate=0.1, max_depth=8",
        "Neural Network (MLP): 256 -> 128 -> 64 neurons, ReLU, Adam optimiser",
        "Ensemble: Average probability vectors of all 3 models -> final prediction",
        "Saved: .pkl files for all models + scaler + encoder + metrics.json",
    ], title_color=GREEN)

    # SLIDE 6: ATTACK TYPES
    slide = add_slide(prs)
    set_bg(slide, 13, 17, 31)
    add_rect(slide, 0, 0, 10, 0.08, 0, 200, 230)
    add_textbox(slide, "Attack Types Detected (7 Classes)", 0.3, 0.15, 9.2, 0.6,
                size=22, bold=True, color=CYAN)
    add_rect(slide, 0.3, 0.82, 8.5, 0.03, 0, 160, 200)
    attacks_ppt = [
        ("Normal Traffic", "NONE",     (0,  200, 100), "Low count, balanced bytes, no errors"),
        ("DDoS",           "CRITICAL", (220, 30,  60), "count=400-500, SYN rate~1.0, ICMP"),
        ("Port Scan",      "MEDIUM",   (200,130,   0), "High count, zero dst bytes, many ports"),
        ("Brute Force",    "HIGH",     (220, 80,   0), "failed_logins=10-20, TCP"),
        ("SQL Injection",  "HIGH",     (220, 80,   0), "hot indicators=15-30, large bytes"),
        ("XSS",            "MEDIUM",   (200,130,   0), "Moderate hot, shorter than SQL"),
        ("Botnet C&C",     "CRITICAL", (220, 30,  60), "duration=100-300s, UDP, steady flow"),
    ]
    y = 1.05
    for name, sev, col, sig in attacks_ppt:
        add_rect(slide, 0.3, y, 9.0, 0.72, 0, 30, 60)
        add_textbox(slide, name, 0.45, y+0.08, 2.2, 0.5, size=12, bold=True,
                    color=RGBColor(*col))
        add_textbox(slide, sev, 2.7, y+0.08, 1.4, 0.5, size=11, bold=True,
                    color=RGBColor(*col))
        add_textbox(slide, sig, 4.15, y+0.08, 5.0, 0.5, size=11, color=LGREY)
        y += 0.78

    # SLIDE 7: MODEL PERFORMANCE
    slide = add_slide(prs)
    set_bg(slide, 13, 17, 31)
    add_rect(slide, 0, 0, 10, 0.08, 0, 200, 230)
    add_textbox(slide, "Model Performance Results", 0.3, 0.15, 9, 0.6, size=22, bold=True, color=CYAN)
    add_rect(slide, 0.3, 0.82, 8.5, 0.03, 0, 160, 200)
    add_rect(slide, 0.3, 1.0, 9.0, 0.5, 0, 200, 230)
    for i, h in enumerate(["Model", "Accuracy", "Precision", "Recall", "F1-Score"]):
        add_textbox(slide, h, 0.3 + i*1.8, 1.05, 1.75, 0.4,
                    size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    model_rows = [
        ("Random Forest",    "99.90%", "99.90%", "99.90%", "99.90%", (0, 210, 100)),
        ("XGBoost",          "99.87%", "99.87%", "99.87%", "99.87%", (0, 200, 230)),
        ("Neural Net (MLP)", "99.64%", "99.65%", "99.64%", "99.64%", (124, 77, 255)),
    ]
    for ri, (m, acc, pr, rec, f1, col) in enumerate(model_rows):
        yy = 1.55 + ri * 0.75
        bg = (0, 35, 60) if ri % 2 == 0 else (0, 25, 50)
        add_rect(slide, 0.3, yy, 9.0, 0.7, *bg)
        add_textbox(slide, m, 0.35, yy+0.1, 1.7, 0.5, size=12, bold=True,
                    color=RGBColor(*col))
        for ci, val in enumerate([acc, pr, rec, f1], 1):
            add_textbox(slide, val, 0.3+ci*1.8, yy+0.1, 1.7, 0.5, size=14, bold=True,
                        color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide,
        "All models evaluated on 20,000-record held-out test set | Ensemble reduces error further",
        0.3, 4.0, 9.2, 0.5, size=11, color=LGREY, align=PP_ALIGN.CENTER, italic=True)

    # SLIDE 8: FRONTEND FEATURES
    bullet_slide(prs, "Frontend Dashboard Features", [
        "SYSTEM ONLINE/OFFLINE indicator - health-checks /api/status on load",
        "Live Threat Feed - scrolling feed updated every 1.8s via /api/simulate",
        "Traffic Timeline - Line chart showing normal vs threat volume over time",
        "Threat Distribution - Doughnut chart with real-time attack type percentages",
        "Model Performance Bars - animated accuracy bars from /api/status metrics",
        "Packet Analyser - Submit custom packet features, choose model, see instant results",
        "Quick Presets - One-click fill for Normal, DDoS, Port Scan, Brute Force, SQL, Botnet",
        "Threat Log Table - Tabular history of last 40 classified packets from /api/history",
    ], title_color=CYAN)

    # SLIDE 9: DATA FLOW
    slide = add_slide(prs)
    set_bg(slide, 13, 17, 31)
    add_rect(slide, 0, 0, 10, 0.08, 0, 200, 230)
    add_textbox(slide, "How It Works - Data Flow", 0.3, 0.15, 9, 0.6, size=22, bold=True, color=CYAN)
    add_rect(slide, 0.3, 0.82, 8.5, 0.03, 0, 160, 200)
    flow_steps = [
        ("1", "Browser loads index.html, JS boot() fires", CYAN),
        ("2", "GET /api/status -> sets ONLINE indicator + loads model metrics", GREEN),
        ("3", "startScan() polls GET /api/simulate every 1800ms", YELLOW),
        ("4", "Flask randomly picks attack type, generates 20 features", LGREY),
        ("5", "3 ML models each produce probability vector (7 classes)", LGREY),
        ("6", "Probabilities averaged -> final threat_type + severity + confidence", GREEN),
        ("7", "JSON response updates Feed, Donut chart, Stats, Alert banner", CYAN),
        ("8", "Manual: POST /api/predict -> single packet full analysis", YELLOW),
    ]
    y = 1.05
    for num, text, col in flow_steps:
        add_rect(slide, 0.3, y, 0.55, 0.58, 0, 160, 200)
        add_textbox(slide, num, 0.3, y+0.05, 0.55, 0.5, size=16, bold=True,
                    color=DARK, align=PP_ALIGN.CENTER)
        add_rect(slide, 0.92, y, 8.7, 0.58, 0, 25, 50)
        add_textbox(slide, text, 1.05, y+0.08, 8.4, 0.45, size=12, color=col)
        y += 0.68

    # SLIDE 10: FUTURE WORK & CONCLUSION
    slide = add_slide(prs)
    set_bg(slide, 13, 17, 31)
    add_rect(slide, 0, 0, 10, 0.08, 0, 200, 230)
    add_textbox(slide, "Future Work and Conclusion", 0.3, 0.15, 9, 0.6, size=22, bold=True, color=CYAN)
    add_rect(slide, 0.3, 0.82, 8.5, 0.03, 0, 160, 200)
    add_rect(slide, 0.3, 1.0, 4.4, 5.6, 0, 30, 60)
    add_textbox(slide, "Future Enhancements", 0.5, 1.1, 3.9, 0.4, size=14, bold=True, color=YELLOW)
    future_items = [
        "Real packet capture (Scapy/libpcap)",
        "LSTM deep learning for time-series",
        "SQLite persistent threat logging",
        "Email/SMS critical alerts (Twilio)",
        "Docker containerised deployment",
        "JWT authentication for dashboard",
        "VirusTotal IP reputation lookup",
    ]
    yy = 1.6
    for fi in future_items:
        add_textbox(slide, "> " + fi, 0.5, yy, 4.0, 0.42, size=11, color=WHITE)
        yy += 0.44
    add_rect(slide, 5.0, 1.0, 4.7, 5.6, 0, 30, 60)
    add_textbox(slide, "Conclusion", 5.2, 1.1, 4.2, 0.4, size=14, bold=True, color=GREEN)
    conc_items = [
        "End-to-end AI threat detection system",
        "3 ML models: 99.6-99.9% accuracy",
        "Real-time dashboard with live charts",
        "Ensemble for robust classification",
        "7 attack types with severity ratings",
        "Full demo ready for examiner",
        "Extensible for real-world deployment",
    ]
    yy = 1.6
    for ci in conc_items:
        add_textbox(slide, "> " + ci, 5.2, yy, 4.3, 0.42, size=11, color=WHITE)
        yy += 0.44

    out_path = os.path.join(OUTPUT_DIR, "CyberShield_Presentation.pptx")
    prs.save(out_path)
    print("[OK] PowerPoint saved: " + out_path)
    return out_path


# ==============================================================
#  ZIP
# ==============================================================
def build_zip(files):
    zip_path = os.path.join(OUTPUT_DIR, "CyberShield_Submission.zip")
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for f in files:
            zf.write(f, os.path.basename(f))
    print("[OK] ZIP saved: " + zip_path)
    return zip_path


# ==============================================================
#  MAIN
# ==============================================================
if __name__ == "__main__":
    print("\n" + "="*55)
    print("  CyberShield Documentation Generator")
    print("="*55 + "\n")

    pdf_path      = build_pdf()
    abstract_path = build_abstract()
    ppt_path      = build_ppt()
    zip_path      = build_zip([pdf_path, abstract_path, ppt_path])

    print("\n" + "="*55)
    print("  All files saved to: docs_output/")
    print("  - CyberShield_Full_Report.pdf")
    print("  - CyberShield_Abstract.pdf")
    print("  - CyberShield_Presentation.pptx")
    print("  - CyberShield_Submission.zip")
    print("="*55 + "\n")
